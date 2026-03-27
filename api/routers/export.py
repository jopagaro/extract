"""
Export router — download report outputs as formatted files.

Endpoints:
  GET /projects/{project_id}/reports/{run_id}/export?format=json
  GET /projects/{project_id}/reports/{run_id}/export?format=md
  GET /projects/{project_id}/reports/{run_id}/export?format=txt
  GET /projects/{project_id}/reports/{run_id}/export?format=pdf
"""

from __future__ import annotations

import json
import re
from datetime import datetime, timezone
from typing import Literal

from fastapi import APIRouter, HTTPException, Query
from fastapi.responses import Response

from engine.core.paths import project_root, run_root

router = APIRouter(tags=["export"])


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _project_exists(project_id: str) -> None:
    if not project_root(project_id).exists():
        raise HTTPException(status_code=404, detail=f"Project '{project_id}' not found")


def _get_all_sections(project_id: str, run_id: str) -> dict[str, object]:
    rdir = run_root(project_id, run_id)
    if not rdir.exists():
        return {}
    sections = {}
    for jf in sorted(rdir.glob("*.json")):
        if jf.name == "run_status.json":
            continue
        section_name = jf.stem.replace(f"{run_id}_", "")
        with jf.open() as f:
            sections[section_name] = json.load(f)
    return sections


SECTION_TITLES = {
    "07_assembly":            "Analyst Narrative",
    "02_market_intelligence": "Market Intelligence",
    "01_project_facts":       "Project Facts",
    "03_geology":             "Geology & Resources",
    "04_economics":           "Economics & Financial Analysis",
    "05_risks":               "Risks & Uncertainties",
    "06_dcf_model":           "DCF Financial Model",
    "14_metallurgy":          "Metallurgy & Recovery",
    "15_permitting":          "Permitting & Regulatory Status",
    "16_operator":            "Operator & Management Track Record",
    "17_capital_structure":   "Capital Structure & Encumbrances",
    "18_cad_semantics":       "Mine Design (CAD)",
    "08_data_gaps":           "Data Gap Report",
    "09_confidence":          "Confidence Assessment",
    "10_contradictions":      "Contradiction & Consistency Check",
    "12_jurisdiction_risk":   "Jurisdiction Profile",
    "13_compliance":          "NI 43-101 / JORC Compliance Check",
    "11_citations":           "Appendix B -- Source Citations",
    "00_data_sources":        "Appendix A -- Source Documents",
    "project_facts":          "Project Facts",
    "geology_summary":        "Geology",
    "economics_summary":      "Economics",
    "risk_summary":           "Risks",
}


def _format_section(content: object, indent: int = 0) -> str:
    pad = "  " * indent
    if isinstance(content, dict):
        parts = []
        for k, v in content.items():
            label = k.replace("_", " ").title()
            if isinstance(v, (dict, list)):
                parts.append(f"{pad}**{label}:**")
                parts.append(_format_section(v, indent + 1))
            else:
                parts.append(f"{pad}**{label}:** {v}")
        return "\n".join(parts)
    elif isinstance(content, list):
        if not content:
            return f"{pad}_(none)_"
        parts = []
        for item in content:
            if isinstance(item, dict):
                parts.append(_format_section(item, indent))
                parts.append("")
            else:
                parts.append(f"{pad}- {item}")
        return "\n".join(parts)
    else:
        return f"{pad}{content}"


def _sections_to_markdown(project_id: str, run_id: str, sections: dict) -> str:
    lines = [
        "# Extract — Technical Analysis Report",
        "",
        f"**Project:** {project_id}  ",
        f"**Run ID:** {run_id}  ",
        f"**Generated:** {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}",
        "",
        "---",
        "",
    ]
    for section_key, content in sections.items():
        title = SECTION_TITLES.get(section_key, section_key.replace("_", " ").title())
        lines.append(f"## {title}")
        lines.append("")
        lines.append(_format_section(content))
        lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


def _flatten_for_pdf(content: object, lines: list[str], indent: int = 0) -> None:
    pad = "  " * indent
    if isinstance(content, dict):
        for k, v in content.items():
            label = k.replace("_", " ").title()
            if isinstance(v, str) and len(v) > 80:
                lines.append(v)
                lines.append("")
            elif isinstance(v, list):
                if v:
                    lines.append(f"{pad}{label}:")
                    _flatten_for_pdf(v, lines, indent + 1)
                    lines.append("")
            elif isinstance(v, dict):
                _flatten_for_pdf(v, lines, indent)
            elif v is not None and str(v).strip():
                lines.append(f"{pad}{label}: {v}")
    elif isinstance(content, list):
        for item in content:
            if isinstance(item, dict):
                _flatten_for_pdf(item, lines, indent)
                lines.append("")
            else:
                lines.append(f"{pad}* {item}")
    else:
        if str(content).strip():
            lines.append(f"{pad}{content}")


def _safe(text: str) -> str:
    replacements = {
        "\u2014": "-", "\u2013": "-", "\u2012": "-", "\u2015": "-",
        "\u2018": "'", "\u2019": "'", "\u201a": "'",
        "\u201c": '"', "\u201d": '"', "\u201e": '"',
        "\u2022": "-", "\u2023": "-", "\u25cf": "-",
        "\u2026": "...",
        "\u00b0": "deg", "\u00b2": "2", "\u00b3": "3",
        "\u00d7": "x",  "\u00f7": "/", "\u2248": "~",
        "\u2264": "<=", "\u2265": ">=",
        "\u00b1": "+/-",
        "\u20ac": "EUR", "\u00a3": "GBP", "\u00a5": "JPY",
        "\u00ae": "(R)", "\u00a9": "(C)", "\u2122": "(TM)",
        "\u00a0": " ",
    }
    for char, replacement in replacements.items():
        text = text.replace(char, replacement)
    return text.encode("latin-1", errors="replace").decode("latin-1")


# ---------------------------------------------------------------------------
# Design tokens  — Substack-style: off-white, dark ink, one accent bar
# ---------------------------------------------------------------------------

_BG         = (247, 244, 238)   # warm off-white page background
_INK        = (26, 23, 19)      # near-black body text
_INK_SOFT   = (80, 75, 68)      # secondary text / captions
_GRAY       = (150, 144, 136)   # labels, metadata, rules
_RULE_CLR   = (210, 206, 198)   # horizontal dividers
_ACCENT     = (44, 74, 62)      # dark muted green — top bar only
_WHITE      = (255, 255, 255)

# Layout
_LM  = 22   # left margin
_RM  = 22   # right margin
_PW  = 210 - _LM - _RM   # usable width on A4


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _bg(pdf) -> None:
    """Fill the current page with the off-white background."""
    pdf.set_fill_color(*_BG)
    pdf.rect(0, 0, 210, 297, "F")


def _accent_bar(pdf) -> None:
    """Thin colored stripe across the very top of the page."""
    pdf.set_fill_color(*_ACCENT)
    pdf.rect(0, 0, 210, 4, "F")


def _rule(pdf, y: float | None = None) -> None:
    lw_prev = pdf.line_width
    pdf.set_draw_color(*_RULE_CLR)
    pdf.set_line_width(0.25)
    y = y if y is not None else pdf.get_y()
    pdf.line(_LM, y, _LM + _PW, y)
    pdf.set_line_width(lw_prev)


def _label(pdf, text: str) -> None:
    """Small all-caps gray label — used above section titles."""
    pdf.set_font("Helvetica", "B", 7.5)
    pdf.set_text_color(*_GRAY)
    pdf.cell(0, 5, _safe(text.upper()), ln=True)


def _h1(pdf, text: str) -> None:
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(*_INK)
    pdf.multi_cell(_PW, 9, _safe(text))
    pdf.ln(1)


def _h2(pdf, text: str) -> None:
    pdf.set_font("Helvetica", "B", 14)
    pdf.set_text_color(*_INK)
    pdf.multi_cell(_PW, 7, _safe(text))
    pdf.ln(2)


def _body(pdf, text: str, size: float = 10.5) -> None:
    pdf.set_font("Helvetica", "", size)
    pdf.set_text_color(*_INK)
    pdf.multi_cell(_PW, 6.2, _safe(text))


_FIG_RE = re.compile(r'\{\{FIGURE:\s*([^|]+)\|([^}]+)\}\}')


def _prose(pdf, text: str, size: float = 10.5,
           omf_renders_dir=None) -> None:
    """Render multi-paragraph prose, substituting {{FIGURE: ...}} placeholders."""

    pdf.set_font("Helvetica", "", size)
    pdf.set_text_color(*_INK)

    # Split on figure placeholders first, then split remaining text on blank lines
    segments = _FIG_RE.split(text)
    # _FIG_RE.split produces: [text, filename, caption, text, filename, caption, ...]
    i = 0
    while i < len(segments):
        if i % 3 == 0:
            # Plain prose segment
            for para in [p.strip() for p in segments[i].split("\n\n") if p.strip()]:
                pdf.set_font("Helvetica", "", size)
                pdf.set_text_color(*_INK)
                pdf.multi_cell(_PW, 6.2, _safe(" ".join(para.split())))
                pdf.ln(4)
        elif i % 3 == 1:
            # Figure: segments[i] = filename, segments[i+1] = caption
            filename = segments[i].strip()
            caption  = segments[i + 1].strip() if i + 1 < len(segments) else filename
            # Look for render in normalized/renders/ first, then raw/renders/
            img_path = None
            if omf_renders_dir is not None:
                candidate = omf_renders_dir / filename
                if candidate.exists():
                    img_path = candidate
            if img_path is None:
                # Try to find in project renders dirs (passed via closure in _generate_pdf)
                pass
            if img_path and img_path.suffix.lower() in {".png", ".jpg", ".jpeg"}:
                try:
                    pdf.ln(4)
                    pdf.image(str(img_path), w=_PW)
                    pdf.ln(2)
                    _caption(pdf, caption)
                    pdf.ln(6)
                except Exception:
                    _caption(pdf, f"[Figure not available: {filename}]")
            else:
                _caption(pdf, f"[Figure: {filename} — {caption}]")
            i += 1  # skip caption segment (already consumed)
        i += 1


def _caption(pdf, text: str) -> None:
    pdf.set_font("Helvetica", "I", 8.5)
    pdf.set_text_color(*_GRAY)
    pdf.multi_cell(_PW, 5, _safe(text))


def _section_break(pdf, title: str, subtitle: str | None = None, num: str | None = None) -> None:
    """Clean section opener: optional number label, bold title, optional subtitle, rule."""
    pdf.ln(2)
    if num:
        _label(pdf, f"Section {num}")
        pdf.ln(1)
    _h2(pdf, title)
    if subtitle:
        _caption(pdf, subtitle)
        pdf.ln(1)
    _rule(pdf)
    pdf.ln(6)


# ---------------------------------------------------------------------------
# Data renderers
# ---------------------------------------------------------------------------

def _kv_list(pdf, items: list[tuple[str, str]]) -> None:
    """Render a clean label: value list, no backgrounds."""
    for label, value in items:
        if not value or str(value).strip() in ("", "None", "null"):
            continue
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(*_INK_SOFT)
        pdf.write(6, _safe(label + ":  "))
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*_INK)
        pdf.write(6, _safe(str(value)))
        pdf.ln(6)


def _simple_table(pdf, headers: list[str], rows: list[list[str]], col_widths: list[float] | None = None) -> None:
    """Minimal table: thin gray borders, no background fills."""
    if not rows:
        return
    if col_widths is None:
        col_widths = [_PW / len(headers)] * len(headers)
    row_h = 6.5

    # Header row
    pdf.set_font("Helvetica", "B", 8.5)
    pdf.set_text_color(*_INK_SOFT)
    pdf.set_draw_color(*_RULE_CLR)
    pdf.set_line_width(0.25)
    x0 = _LM
    for i, h in enumerate(headers):
        pdf.set_xy(x0, pdf.get_y())
        pdf.cell(col_widths[i], row_h, _safe(h), border="B")
        x0 += col_widths[i]
    pdf.ln(row_h)

    # Body
    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*_INK)
    for row in rows:
        if pdf.get_y() > 265:
            pdf.add_page()
        x0 = _LM
        for i, cell in enumerate(row):
            pdf.set_xy(x0, pdf.get_y())
            pdf.cell(col_widths[i], row_h, _safe(str(cell)), border=0)
            x0 += col_widths[i]
        pdf.ln(row_h)
        # thin rule between rows
        _rule(pdf)

    pdf.ln(6)


# ---------------------------------------------------------------------------
# Section order
# ---------------------------------------------------------------------------

_PDF_SECTION_ORDER = [
    "07_assembly",
    "02_market_intelligence",
    "03_geology",
    "04_economics",
    "05_risks",
    "06_dcf_model",
    "14_metallurgy",
    "15_permitting",
    "16_operator",
    "17_capital_structure",
    "18_cad_semantics",
    "08_data_gaps",
    "09_confidence",
    "10_contradictions",
    "12_jurisdiction_risk",
    "13_compliance",
    "00_data_sources",
    "11_citations",
    "01_project_facts",
]

_PDF_SECTION_META: dict[str, dict] = {
    "07_assembly":            {"title": "Analyst Narrative",                        "subtitle": None,                                                                              "num": None},
    "02_market_intelligence": {"title": "Market Intelligence",                      "subtitle": "Live prices and current market context",                                          "num": None},
    "03_geology":             {"title": "Geology & Resources",                      "subtitle": "Deposit geology and resource assessment",                                         "num": "1"},
    "04_economics":           {"title": "Economics & Financial Analysis",            "subtitle": "Capital costs, operating costs, projections",                                     "num": "2"},
    "05_risks":               {"title": "Risks & Uncertainties",                    "subtitle": "Material risks and mitigations",                                                  "num": "3"},
    "06_dcf_model":           {"title": "DCF Financial Model",                      "subtitle": "Discounted cash flow analysis",                                                   "num": "4"},
    "14_metallurgy":          {"title": "Metallurgy & Recovery",                    "subtitle": "Process route, testwork results, and recovery assumptions that drive revenue",    "num": "5"},
    "15_permitting":          {"title": "Permitting & Regulatory Status",           "subtitle": "Permit status, environmental assessment, water rights, Indigenous consultation",   "num": "6"},
    "16_operator":            {"title": "Operator & Management Track Record",       "subtitle": "Management team, prior project outcomes, and company delivery history",           "num": "7"},
    "17_capital_structure":   {"title": "Capital Structure & Encumbrances",         "subtitle": "Shares, warrants, streaming deals, royalty stacks, debt, and funding status",    "num": "8"},
    "18_cad_semantics":       {"title": "Mine Design (CAD)",                        "subtitle": "Structured parameters extracted from CAD and 3D model files",                    "num": "9"},
    "08_data_gaps":           {"title": "Data Gap Report",                          "subtitle": "Material information gaps and recommended actions",                               "num": None},
    "09_confidence":          {"title": "Confidence Assessment",                    "subtitle": "How much trust to place in each section of this report",                         "num": None},
    "10_contradictions":      {"title": "Contradiction & Consistency Check",        "subtitle": "Internal contradictions, numeric mismatches, arithmetic errors",                  "num": None},
    "12_jurisdiction_risk":   {"title": "Jurisdiction Profile",                     "subtitle": "Current mining fiscal regime, royalty structure, and regulatory environment",     "num": None},
    "13_compliance":          {"title": "NI 43-101 / JORC Compliance Check",        "subtitle": "Assessment against NI 43-101 and JORC Code 2012 requirements",                   "num": None},
    "00_data_sources":        {"title": "Appendix A -- Source Documents",           "subtitle": None,                                                                              "num": None},
    "11_citations":           {"title": "Appendix B -- Source Citations",           "subtitle": "Traceability index mapping report claims to source documents",                    "num": None},
    "01_project_facts":       {"title": "Project Facts",                            "subtitle": None,                                                                              "num": None},
}


# ---------------------------------------------------------------------------
# PDF generator
# ---------------------------------------------------------------------------

def _generate_pdf(project_id: str, run_id: str, sections: dict) -> bytes:
    import re
    from fpdf import FPDF

    raw_dir      = project_root(project_id) / "raw" / "documents"
    renders_dir  = project_root(project_id) / "raw" / "renders"
    omf_renders  = project_root(project_id) / "normalized" / "renders"
    IMAGE_EXTS   = {".png", ".jpg", ".jpeg", ".tiff"}

    _FIGURE_RE = re.compile(r'\{\{FIGURE:\s*([^|]+)\|([^}]+)\}\}')

    class ReportPDF(FPDF):
        def header(self):
            _bg(self)
            _accent_bar(self)
            if self.page_no() == 1:
                return
            # Subtle running header
            self.set_y(10)
            self.set_font("Helvetica", "", 8)
            self.set_text_color(*_GRAY)
            proj = _safe(project_id.replace("_", " ").replace("-", " ").title())
            self.cell(0, 5, proj, align="L")
            self.set_y(10)
            self.cell(0, 5, "Extract", align="R")
            self.set_y(17)
            _rule(self)
            self.ln(5)

        def footer(self):
            self.set_y(-14)
            self.set_font("Helvetica", "", 8)
            self.set_text_color(*_GRAY)
            self.cell(0, 6, "For internal research purposes only. Not investment advice.", align="L")
            self.set_y(-14)
            self.cell(0, 6, f"{self.page_no() - 1}", align="R")

    pdf = ReportPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_margins(_LM, 18, _RM)

    # ── Cover ────────────────────────────────────────────────────────────────
    pdf.add_page()

    pdf.set_y(22)

    # Publication label
    _label(pdf, "Extract  —  Technical Analysis Report")
    pdf.ln(5)

    # Project title
    proj_display = project_id.replace("_", " ").replace("-", " ").title()
    _h1(pdf, proj_display)
    pdf.ln(3)

    # Study level / stage
    assembly = sections.get("07_assembly", {})
    stage = ""
    if isinstance(assembly, dict):
        stage = assembly.get("project_stage") or assembly.get("study_level") or ""
    if stage and stage.lower() not in ("unknown", "not specified", ""):
        pdf.set_font("Helvetica", "", 12)
        pdf.set_text_color(*_INK_SOFT)
        pdf.cell(0, 7, _safe(str(stage)), ln=True)
        pdf.ln(2)

    _rule(pdf)
    pdf.ln(6)

    # Meta line
    date_str = datetime.now(timezone.utc).strftime("%B %d, %Y")
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(*_GRAY)
    pdf.cell(0, 5, _safe(f"{date_str}  -  Run {run_id}  -  Internal / Confidential"), ln=True)
    pdf.ln(10)

    # Disclaimer
    _caption(pdf,
        "This report is generated by an AI system for internal research purposes only. "
        "It does not constitute investment advice or a formal technical study. "
        "All figures and conclusions should be verified against primary source documents."
    )

    # ── Content pages ────────────────────────────────────────────────────────
    ordered_keys = _PDF_SECTION_ORDER + [k for k in sections if k not in _PDF_SECTION_ORDER]

    for section_key in ordered_keys:
        if section_key not in sections:
            continue
        content  = sections[section_key]
        meta_s   = _PDF_SECTION_META.get(section_key, {})
        title    = meta_s.get("title") or SECTION_TITLES.get(section_key, section_key.replace("_", " ").title())
        subtitle = meta_s.get("subtitle")
        num      = meta_s.get("num")

        pdf.add_page()
        _section_break(pdf, title, subtitle, num)

        # ── Analyst Narrative ─────────────────────────────────────────────
        if section_key == "07_assembly" and isinstance(content, dict):
            # Key callouts — plain inline list, no boxes
            callouts = content.get("key_callouts")
            if isinstance(callouts, list) and callouts:
                _label(pdf, "Key metrics")
                pdf.ln(2)
                for c in callouts:
                    v = _safe(str(c.get("value", "")))
                    l = _safe(str(c.get("label", "")))
                    ctx = c.get("context", "")
                    pdf.set_font("Helvetica", "B", 10.5)
                    pdf.set_text_color(*_INK)
                    pdf.write(6.5, v + "  ")
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(*_INK_SOFT)
                    pdf.write(6.5, _safe(l + ("  -  " + str(ctx) if ctx else "")))
                    pdf.ln(7)
                pdf.ln(4)
                _rule(pdf)
                pdf.ln(6)

            # Narrative prose
            narrative = content.get("narrative")
            if isinstance(narrative, str):
                _prose(pdf, narrative)

            # Analyst conclusion
            conclusion = content.get("analyst_conclusion")
            if isinstance(conclusion, str):
                pdf.ln(4)
                _rule(pdf)
                pdf.ln(5)
                _label(pdf, "Analyst conclusion")
                pdf.ln(2)
                pdf.set_font("Helvetica", "I", 10.5)
                pdf.set_text_color(*_INK)
                pdf.multi_cell(_PW, 6.2, _safe(conclusion))
                pdf.ln(2)

            # Consistency flags — plain text, no box
            flags = content.get("consistency_flags")
            if isinstance(flags, list) and flags:
                pdf.ln(4)
                _rule(pdf)
                pdf.ln(5)
                _label(pdf, "Consistency notes")
                pdf.ln(2)
                for flag in flags:
                    pdf.set_font("Helvetica", "", 9.5)
                    pdf.set_text_color(*_INK_SOFT)
                    pdf.multi_cell(_PW, 5.8, _safe(f"- {flag}"))
                    pdf.ln(1)

        # ── Market Intelligence ───────────────────────────────────────────
        elif section_key == "02_market_intelligence" and isinstance(content, dict):
            if content.get("error"):
                _caption(pdf, str(content["error"]))
            else:
                gathered_at = content.get("gathered_at", "")
                if gathered_at:
                    _caption(pdf, f"Data gathered: {gathered_at}")
                    pdf.ln(4)

                # Live prices — inline label: value list
                live_prices = content.get("live_prices", {}).get("prices", {})
                macro_inds  = content.get("macro_indicators", {}).get("indicators", {})
                price_items: list[tuple[str, str]] = []
                for name, data in live_prices.items():
                    v = data.get("price")
                    u = data.get("unit", "")
                    if v is not None:
                        val_str = f"{v:,.2f} {u}".strip() if isinstance(v, (int, float)) else f"{v} {u}".strip()
                        price_items.append((name.replace("_", " ").title(), val_str))
                for name, data in macro_inds.items():
                    v = data.get("value")
                    u = data.get("unit", "")
                    if v is not None:
                        val_str = f"{v:,.2f} {u}".strip() if isinstance(v, (int, float)) else f"{v} {u}".strip()
                        price_items.append((name.replace("_", " ").title(), val_str))

                if price_items:
                    _label(pdf, "Live market data")
                    pdf.ln(2)
                    _kv_list(pdf, price_items)
                    pdf.ln(4)
                    _rule(pdf)
                    pdf.ln(6)

                # Project intelligence
                proj_intel = content.get("project_intelligence", {})
                if proj_intel.get("findings"):
                    subject = proj_intel.get("search_subject", "Project")
                    _label(pdf, f"Project intelligence — {subject}")
                    pdf.ln(2)
                    _prose(pdf, proj_intel["findings"], size=10)

                # Commodity market
                comm = content.get("commodity_market", {})
                if comm.get("analysis"):
                    _label(pdf, f"Commodity market — {comm.get('commodity', '')}")
                    pdf.ln(2)
                    _prose(pdf, comm["analysis"], size=10)

                # Macro context
                macro = content.get("macro_context", {})
                if macro.get("analysis"):
                    _label(pdf, "Macroeconomic context")
                    pdf.ln(2)
                    _prose(pdf, macro["analysis"], size=10)

                notice = content.get("notice")
                if notice:
                    pdf.ln(2)
                    _caption(pdf, notice)

        # ── DCF Financial Model ───────────────────────────────────────────
        elif section_key == "06_dcf_model" and isinstance(content, dict):
            if not content.get("model_ran"):
                reason = content.get("reason", "DCF model did not run.")
                _caption(pdf, str(reason))
            else:
                notes = content.get("assumptions_notes")
                if isinstance(notes, str):
                    _caption(pdf, notes)
                    pdf.ln(5)

                summary = content.get("summary")
                if isinstance(summary, dict):
                    _label(pdf, "Valuation summary")
                    pdf.ln(2)
                    skip = {"project_id", "scenario", "after_tax", "notes", "aisc_unit"}
                    kv = [
                        (k.replace("_", " ").title(),
                         f"{v:,}" if isinstance(v, (int, float)) else str(v))
                        for k, v in summary.items()
                        if k not in skip and v is not None
                    ]
                    _kv_list(pdf, kv)
                    pdf.ln(4)
                    _rule(pdf)
                    pdf.ln(6)

                cash_flows = content.get("cash_flows")
                if isinstance(cash_flows, list) and cash_flows:
                    _label(pdf, "Annual cash flows")
                    pdf.ln(3)

                    def _fmt(v):
                        if v is None: return "-"
                        if isinstance(v, (int, float)): return f"{v:,.0f}"
                        return str(v)

                    headers   = ["Year", "Revenue", "OpEx", "CapEx", "Tax", "FCF", "PV"]
                    col_widths = [16, 28, 26, 26, 22, 28, 28]
                    rows_data  = [
                        [
                            str(cf.get("year", "")),
                            _fmt(cf.get("gross_revenue")),
                            _fmt(cf.get("opex")),
                            _fmt(cf.get("capex")),
                            _fmt(cf.get("tax")),
                            _fmt(cf.get("free_cash_flow")),
                            _fmt(cf.get("present_value")),
                        ]
                        for cf in cash_flows[:20]
                    ]
                    _simple_table(pdf, headers, rows_data, col_widths)

        # ── Data sources ──────────────────────────────────────────────────
        elif section_key == "00_data_sources" and isinstance(content, dict):
            notice = content.get("notice")
            if isinstance(notice, str):
                _caption(pdf, notice)
                pdf.ln(6)

            files = content.get("source_files", [])
            if files:
                _label(pdf, f"Source documents ({len(files)})")
                pdf.ln(3)
                for f in files:
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(*_INK)
                    pdf.cell(0, 6.5, _safe(f"  {f}"), ln=True)
                pdf.ln(5)

            # Embedded images
            image_files  = content.get("image_files", [])
            render_files = content.get("render_files", [])
            all_images   = list(dict.fromkeys(image_files + render_files))
            shown = 0
            for img_name in all_images:
                img_path = raw_dir / img_name
                if not img_path.exists():
                    img_path = renders_dir / img_name
                if not img_path.exists() or img_path.suffix.lower() not in IMAGE_EXTS:
                    continue
                if shown == 0:
                    pdf.ln(4)
                    _label(pdf, "Visual references")
                    pdf.ln(3)
                try:
                    pdf.image(str(img_path), w=min(_PW, 130))
                    _caption(pdf, img_name)
                    pdf.ln(4)
                    shown += 1
                except Exception:
                    pass

        # ── Data Gap Report ───────────────────────────────────────────────
        elif section_key == "08_data_gaps" and isinstance(content, dict):
            overall = content.get("overall_data_quality_comment")
            if isinstance(overall, str):
                _prose(pdf, overall)
                pdf.ln(2)
                _rule(pdf)
                pdf.ln(6)

            # Summary counts
            counts = []
            crit = content.get("critical_gaps_count")
            imp  = content.get("important_gaps_count")
            minor = content.get("minor_gaps_count")
            if crit is not None: counts.append(f"{crit} critical")
            if imp  is not None: counts.append(f"{imp} important")
            if minor is not None: counts.append(f"{minor} minor")
            if counts:
                _label(pdf, "Gap summary")
                pdf.ln(2)
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(*_INK_SOFT)
                pdf.cell(0, 6, _safe("  ·  ".join(counts)), ln=True)
                pdf.ln(5)

            gaps = content.get("data_gaps", [])
            if isinstance(gaps, list) and gaps:
                _label(pdf, f"Identified gaps ({len(gaps)})")
                pdf.ln(4)
                for gap in gaps:
                    if not isinstance(gap, dict): continue
                    if pdf.get_y() > 250: pdf.add_page()

                    urgency = str(gap.get("urgency", "")).lower()
                    domain  = str(gap.get("domain", ""))
                    desc    = str(gap.get("gap_description", ""))
                    impact  = str(gap.get("impact_on_analysis", ""))
                    action  = str(gap.get("recommended_action", ""))
                    blocking = gap.get("blocking_advancement", False)

                    # Domain label with urgency tag inline
                    tag = {"critical": "[CRITICAL]", "important": "[IMPORTANT]", "minor": "[MINOR]"}.get(urgency, "")
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_text_color(*_INK)
                    pdf.write(6, _safe(domain))
                    if tag:
                        pdf.set_font("Helvetica", "", 8.5)
                        pdf.set_text_color(*_GRAY)
                        pdf.write(6, f"  {tag}")
                        if blocking:
                            pdf.write(6, "  [blocks advancement]")
                    pdf.ln(6)

                    if desc and desc != "No material gaps identified":
                        pdf.set_font("Helvetica", "", 10)
                        pdf.set_text_color(*_INK)
                        pdf.multi_cell(_PW, 5.8, _safe(desc))
                        pdf.ln(2)

                    if impact:
                        pdf.set_font("Helvetica", "I", 9.5)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.5, _safe(f"Impact: {impact}"))
                        pdf.ln(1)

                    if action:
                        pdf.set_font("Helvetica", "", 9.5)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.5, _safe(f"Action: {action}"))

                    pdf.ln(5)
                    _rule(pdf)
                    pdf.ln(5)

        # ── Confidence Assessment ─────────────────────────────────────────
        elif section_key == "09_confidence" and isinstance(content, dict):
            overall = content.get("overall_confidence_statement")
            if isinstance(overall, str):
                _prose(pdf, overall)
                pdf.ln(2)
                _rule(pdf)
                pdf.ln(6)

            best = content.get("most_reliable_aspect")
            worst = content.get("least_reliable_aspect")
            if best or worst:
                if best:
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_text_color(*_INK_SOFT)
                    pdf.write(6.2, "Most reliable:  ")
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(*_INK)
                    pdf.write(6.2, _safe(str(best)))
                    pdf.ln(7)
                if worst:
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_text_color(*_INK_SOFT)
                    pdf.write(6.2, "Least reliable:  ")
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(*_INK)
                    pdf.write(6.2, _safe(str(worst)))
                    pdf.ln(7)
                pdf.ln(4)
                _rule(pdf)
                pdf.ln(6)

            domains = content.get("domain_confidence", [])
            if isinstance(domains, list) and domains:
                _label(pdf, f"Domain breakdown ({len(domains)} areas)")
                pdf.ln(4)
                for d in domains:
                    if not isinstance(d, dict): continue
                    if pdf.get_y() > 250: pdf.add_page()

                    domain_name = str(d.get("domain", ""))
                    descriptor  = str(d.get("confidence_descriptor", ""))
                    supporting  = d.get("supporting_factors", [])
                    limiting    = d.get("limiting_factors", [])

                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_text_color(*_INK)
                    pdf.multi_cell(_PW, 6, _safe(domain_name))
                    pdf.ln(1)

                    if descriptor:
                        pdf.set_font("Helvetica", "I", 10)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.8, _safe(descriptor))
                        pdf.ln(3)

                    if supporting:
                        pdf.set_font("Helvetica", "", 9)
                        pdf.set_text_color(*_GRAY)
                        pdf.write(5.5, "Supports: ")
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.5, _safe("  ·  ".join(str(f) for f in supporting)))
                        pdf.ln(1)

                    if limiting:
                        pdf.set_font("Helvetica", "", 9)
                        pdf.set_text_color(*_GRAY)
                        pdf.write(5.5, "Limits: ")
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.5, _safe("  ·  ".join(str(f) for f in limiting)))

                    pdf.ln(5)
                    _rule(pdf)
                    pdf.ln(5)

        # ── Contradiction & Consistency Check ────────────────────────────
        elif section_key == "10_contradictions" and isinstance(content, dict):
            overall = content.get("overall_consistency_comment")
            if isinstance(overall, str):
                _prose(pdf, overall)
                pdf.ln(2)
                _rule(pdf)
                pdf.ln(6)

            # Summary counts
            crit_c = content.get("critical_count")
            sig_c  = content.get("significant_count")
            min_c  = content.get("minor_count")
            count_parts = []
            if crit_c is not None: count_parts.append(f"{crit_c} critical")
            if sig_c  is not None: count_parts.append(f"{sig_c} significant")
            if min_c  is not None: count_parts.append(f"{min_c} minor")
            if count_parts:
                _label(pdf, "Contradiction summary")
                pdf.ln(2)
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(*_INK_SOFT)
                pdf.cell(0, 6, _safe("  |  ".join(count_parts)), ln=True)
                pdf.ln(5)

            # Arithmetic errors table
            arith_errors = content.get("arithmetic_errors", [])
            if isinstance(arith_errors, list) and arith_errors:
                _label(pdf, f"Arithmetic errors ({len(arith_errors)})")
                pdf.ln(3)
                headers   = ["Field", "Stated", "Calculated", "Discrepancy"]
                col_widths = [50, 36, 36, 44]
                rows_data  = []
                for e in arith_errors:
                    if not isinstance(e, dict): continue
                    disc = e.get("discrepancy_pct")
                    disc_str = f"{disc:.1f}%" if isinstance(disc, (int, float)) else "-"
                    rows_data.append([
                        str(e.get("field", "")),
                        str(e.get("stated_value", "")),
                        str(e.get("calculated_value", "")),
                        disc_str,
                    ])
                _simple_table(pdf, headers, rows_data, col_widths)
                pdf.ln(4)
                _rule(pdf)
                pdf.ln(6)

            # Individual contradictions
            contradictions = content.get("contradictions", [])
            if isinstance(contradictions, list) and contradictions:
                _label(pdf, f"Identified contradictions ({len(contradictions)})")
                pdf.ln(4)
                _SEVER_TAGS = {
                    "critical":    "[CRITICAL]",
                    "significant": "[SIGNIFICANT]",
                    "minor":       "[MINOR]",
                }
                for c in contradictions:
                    if not isinstance(c, dict): continue
                    if pdf.get_y() > 250: pdf.add_page()

                    sev   = str(c.get("severity", "")).lower()
                    ctype = str(c.get("contradiction_type", "")).replace("_", " ").title()
                    desc  = str(c.get("description", ""))
                    loc_a = str(c.get("location_a", ""))
                    val_a = str(c.get("value_a", ""))
                    loc_b = str(c.get("location_b", ""))
                    val_b = str(c.get("value_b", ""))
                    correct = c.get("correct_value")
                    impact = str(c.get("economic_impact", ""))
                    resolution = str(c.get("resolution_required", ""))

                    tag = _SEVER_TAGS.get(sev, "")
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.set_text_color(*_INK)
                    pdf.write(6, _safe(ctype))
                    if tag:
                        pdf.set_font("Helvetica", "", 8.5)
                        pdf.set_text_color(*_GRAY)
                        pdf.write(6, f"  {tag}")
                        if impact:
                            pdf.write(6, f"  impact: {impact}")
                    pdf.ln(6)

                    if desc:
                        pdf.set_font("Helvetica", "", 10)
                        pdf.set_text_color(*_INK)
                        pdf.multi_cell(_PW, 5.8, _safe(desc))
                        pdf.ln(3)

                    # Value comparison
                    if loc_a or val_a:
                        pdf.set_font("Helvetica", "B", 9)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.write(5.5, _safe(f"{loc_a}: ") if loc_a else "")
                        pdf.set_font("Helvetica", "I", 9)
                        pdf.set_text_color(*_INK)
                        pdf.write(5.5, _safe(val_a))
                        pdf.ln(5.5)
                    if loc_b or val_b:
                        pdf.set_font("Helvetica", "B", 9)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.write(5.5, _safe(f"{loc_b}: ") if loc_b else "")
                        pdf.set_font("Helvetica", "I", 9)
                        pdf.set_text_color(*_INK)
                        pdf.write(5.5, _safe(val_b))
                        pdf.ln(5.5)
                    if correct:
                        pdf.set_font("Helvetica", "B", 9)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.write(5.5, "Correct value: ")
                        pdf.set_font("Helvetica", "", 9)
                        pdf.set_text_color(*_INK)
                        pdf.write(5.5, _safe(str(correct)))
                        pdf.ln(5.5)

                    if resolution:
                        pdf.ln(1)
                        pdf.set_font("Helvetica", "", 9.5)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.5, _safe(f"Resolution: {resolution}"))

                    pdf.ln(5)
                    _rule(pdf)
                    pdf.ln(5)

        # ── Source Citations ──────────────────────────────────────────────
        elif section_key == "11_citations" and isinstance(content, dict):
            comment = content.get("citation_coverage_comment")
            if isinstance(comment, str):
                _prose(pdf, comment)
                pdf.ln(2)
                _rule(pdf)
                pdf.ln(6)

            # Summary counts
            total_c    = content.get("total_citations")
            not_found  = content.get("not_found_count")
            count_parts = []
            if total_c   is not None: count_parts.append(f"{total_c} citations")
            if not_found is not None: count_parts.append(f"{not_found} not found in sources")
            if count_parts:
                _label(pdf, "Citation summary")
                pdf.ln(2)
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(*_INK_SOFT)
                pdf.cell(0, 6, _safe("  |  ".join(count_parts)), ln=True)
                pdf.ln(5)

            uncited = content.get("uncited_sections", [])
            if isinstance(uncited, list) and uncited:
                pdf.set_font("Helvetica", "", 9.5)
                pdf.set_text_color(*_INK_SOFT)
                pdf.multi_cell(_PW, 5.5, _safe("Sections with no citations: " + ", ".join(str(u) for u in uncited)))
                pdf.ln(5)

            citations = content.get("citations", [])
            _SECTION_LABELS = {
                "03_geology":   "Geology",
                "04_economics": "Economics",
                "05_risks":     "Risks",
                "06_dcf_model": "DCF Model",
                "07_assembly":  "Narrative",
            }
            _CONF_TAGS = {
                "direct":    "[DIRECT]",
                "inferred":  "[INFERRED]",
                "not_found": "[NOT FOUND]",
            }
            if isinstance(citations, list) and citations:
                _label(pdf, f"Citation index ({len(citations)})")
                pdf.ln(4)

                current_section = None
                for c in citations:
                    if not isinstance(c, dict): continue
                    if pdf.get_y() > 255: pdf.add_page()

                    sec = str(c.get("section", ""))
                    if sec != current_section:
                        current_section = sec
                        sec_label = _SECTION_LABELS.get(sec, sec.replace("_", " ").title())
                        pdf.ln(3)
                        _label(pdf, sec_label)
                        pdf.ln(2)

                    cid       = str(c.get("citation_id", ""))
                    claim     = str(c.get("claim", ""))
                    src_file  = str(c.get("source_file", ""))
                    src_quote = str(c.get("source_quote", ""))
                    src_loc   = c.get("location_in_source")
                    confidence= str(c.get("confidence", "direct"))
                    conf_tag  = _CONF_TAGS.get(confidence, "")

                    # Citation ID + confidence tag
                    pdf.set_font("Helvetica", "B", 9)
                    pdf.set_text_color(*_GRAY)
                    pdf.write(5.5, _safe(cid + "  "))
                    if conf_tag:
                        pdf.set_font("Helvetica", "", 8.5)
                        pdf.write(5.5, conf_tag)
                    pdf.ln(5.5)

                    # Claim text
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(*_INK)
                    pdf.multi_cell(_PW, 5.8, _safe(claim))
                    pdf.ln(2)

                    # Source info
                    if confidence != "not_found" and src_file:
                        loc_str = f"  —  {src_loc}" if src_loc else ""
                        pdf.set_font("Helvetica", "B", 8.5)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.write(5.2, _safe(src_file + loc_str))
                        pdf.ln(5.5)
                        if src_quote:
                            pdf.set_font("Helvetica", "I", 9)
                            pdf.set_text_color(*_GRAY)
                            pdf.multi_cell(_PW, 5.2, _safe(f'"{src_quote[:200]}"'))
                    elif confidence == "not_found":
                        pdf.set_font("Helvetica", "I", 9)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW, 5.2, _safe("No supporting passage found in source documents."))

                    pdf.ln(5)
                    _rule(pdf)
                    pdf.ln(4)

        # ── NI 43-101 / JORC Compliance ───────────────────────────────────
        elif section_key == "13_compliance" and isinstance(content, dict):
            overall = content.get("overall_status", "")
            std_applied = content.get("standard_applied_for_assessment", "NI 43-101")
            summary_txt = content.get("overall_summary")
            qp = content.get("qualified_person") or {}
            study_type = content.get("study_type")
            study_match = content.get("study_type_resource_match")
            study_match_note = content.get("study_type_resource_match_note")
            checks = content.get("checks") or []
            critical_gaps = content.get("critical_gaps") or []
            minor_gaps = content.get("minor_gaps") or []
            met = content.get("compliant_items_count", 0)
            partial = content.get("partial_items_count", 0)
            missing = content.get("missing_items_count", 0)

            _OVERALL_LABELS = {
                "compliant":          "COMPLIANT",
                "likely_compliant":   "LIKELY COMPLIANT",
                "deficiencies_found": "DEFICIENCIES FOUND",
                "major_gaps":         "MAJOR GAPS IDENTIFIED",
            }
            _OVERALL_COLORS = {
                "compliant":          (22, 101, 52),
                "likely_compliant":   (30, 64, 175),
                "deficiencies_found": (133, 77, 14),
                "major_gaps":         (153, 27, 27),
            }

            # Overall status line
            pdf.set_font("Helvetica", "B", 11)
            pdf.set_text_color(*_OVERALL_COLORS.get(overall, _INK))
            pdf.cell(0, 7, _safe(_OVERALL_LABELS.get(overall, overall.upper())), ln=True)
            pdf.set_text_color(*_INK_SOFT)
            pdf.set_font("Helvetica", "", 9)
            pdf.cell(0, 5, _safe(f"Assessed against {std_applied}  |  Met: {met}   Partial: {partial}   Missing: {missing}"), ln=True)
            pdf.ln(3)
            pdf.set_text_color(*_INK)

            if summary_txt:
                _prose(pdf, str(summary_txt))
                pdf.ln(4)

            # QP / Study type
            if qp or study_type:
                qp_name = qp.get("name") if isinstance(qp, dict) else None
                qp_cred = qp.get("credentials") if isinstance(qp, dict) else None
                qp_named = qp.get("named", False) if isinstance(qp, dict) else False
                _label(pdf, "Qualified / Competent Person")
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(*_INK)
                qp_line = ("Named: " + qp_name) if qp_named and qp_name else ("Not named" if not qp_named else "Not stated")
                if qp_cred:
                    qp_line += f"  ({qp_cred})"
                pdf.cell(0, 6, _safe(qp_line), ln=True)
                pdf.ln(3)

                if study_type:
                    _label(pdf, "Study Type")
                    pdf.set_font("Helvetica", "", 10)
                    match_label = {"ok": "Consistent", "concern": "Concern", "violation": "Violation"}.get(study_match or "", "")
                    st_line = study_type.upper()
                    if match_label:
                        st_line += f"  —  {match_label}"
                    pdf.cell(0, 6, _safe(st_line), ln=True)
                    if study_match_note:
                        pdf.set_font("Helvetica", "I", 9)
                        pdf.set_text_color(*_GRAY)
                        pdf.multi_cell(_PW, 5.2, _safe(str(study_match_note)))
                        pdf.set_text_color(*_INK)
                    pdf.ln(3)

            # Critical gaps
            if critical_gaps:
                _label(pdf, "Critical Gaps")
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(153, 27, 27)
                for gap in critical_gaps:
                    pdf.set_x(_LM + 4)
                    pdf.multi_cell(_PW - 4, 5.5, _safe(f"- {gap}"))
                pdf.set_text_color(*_INK)
                pdf.ln(4)

            # Per-check table
            if checks:
                _label(pdf, "Requirement Checks")
                pdf.ln(2)
                _STATUS_LABELS = {"met": "[MET]", "partial": "[PARTIAL]", "missing": "[MISSING]", "not_applicable": "[N/A]"}
                _STATUS_COLORS = {
                    "met":            (22, 101, 52),
                    "partial":        (133, 77, 14),
                    "missing":        (153, 27, 27),
                    "not_applicable": (107, 114, 128),
                }
                prev_category = None
                for chk in checks:
                    if not isinstance(chk, dict):
                        continue
                    cat  = chk.get("category", "")
                    req  = chk.get("requirement", "")
                    stat = chk.get("status", "not_applicable")
                    finding = chk.get("finding", "")
                    rec  = chk.get("recommendation")

                    if cat and cat != prev_category:
                        pdf.ln(3)
                        pdf.set_font("Helvetica", "B", 9.5)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.cell(0, 6, _safe(cat), ln=True)
                        prev_category = cat

                    # Status badge
                    pdf.set_x(_LM + 2)
                    pdf.set_font("Helvetica", "B", 8.5)
                    pdf.set_text_color(*_STATUS_COLORS.get(stat, _GRAY))
                    pdf.cell(18, 5.5, _safe(_STATUS_LABELS.get(stat, "[?]")))
                    # Requirement
                    pdf.set_font("Helvetica", "", 9.5)
                    pdf.set_text_color(*_INK)
                    pdf.multi_cell(_PW - 18, 5.5, _safe(req))
                    # Finding
                    if finding:
                        pdf.set_x(_LM + 22)
                        pdf.set_font("Helvetica", "", 9)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.multi_cell(_PW - 22, 5.2, _safe(str(finding)[:200]))
                    # Recommendation
                    if rec and stat not in ("met", "not_applicable"):
                        pdf.set_x(_LM + 22)
                        pdf.set_font("Helvetica", "I", 8.5)
                        pdf.set_text_color(133, 77, 14)
                        pdf.multi_cell(_PW - 22, 5, _safe(f"Rec: {str(rec)[:150]}"))
                    pdf.set_text_color(*_INK)
                    pdf.ln(2)

            # Minor gaps
            if minor_gaps:
                pdf.ln(3)
                _label(pdf, "Minor / Technical Gaps")
                pdf.set_font("Helvetica", "", 9.5)
                pdf.set_text_color(133, 77, 14)
                for gap in minor_gaps:
                    pdf.set_x(_LM + 4)
                    pdf.multi_cell(_PW - 4, 5.5, _safe(f"- {gap}"))
                pdf.set_text_color(*_INK)

        # ── Generic sections (geology, economics, risk, project facts) ────
        else:
            flat_lines: list[str] = []
            _flatten_for_pdf(content, flat_lines)

            for raw_line in flat_lines:
                raw_str = str(raw_line)
                if not raw_str.strip():
                    pdf.ln(2)
                    continue

                # Long prose lines may contain {{FIGURE: ...}} placeholders
                if len(raw_str) > 80 and "{{FIGURE:" in raw_str:
                    _prose(pdf, raw_str, omf_renders_dir=omf_renders)
                    continue

                line = _safe(raw_str)
                line = " ".join(w[:120] if len(w) > 120 else w for w in line.split())

                if line.startswith("*"):
                    pdf.set_font("Helvetica", "", 10)
                    pdf.set_text_color(*_INK)
                    pdf.set_x(_LM + 4)
                    pdf.multi_cell(_PW - 4, 6, line)
                    pdf.set_x(_LM)
                elif line.endswith(":") and len(line) < 50:
                    pdf.ln(3)
                    _label(pdf, line.rstrip(":"))
                    pdf.ln(1)
                else:
                    colon_pos = line.find(":")
                    if 0 < colon_pos < 35 and "\n" not in line[:colon_pos]:
                        pdf.set_font("Helvetica", "B", 10)
                        pdf.set_text_color(*_INK_SOFT)
                        pdf.write(6.2, line[:colon_pos + 1] + "  ")
                        pdf.set_font("Helvetica", "", 10)
                        pdf.set_text_color(*_INK)
                        pdf.write(6.2, line[colon_pos + 1:].strip())
                        pdf.ln(6.2)
                    else:
                        pdf.set_font("Helvetica", "", 10.5)
                        pdf.set_text_color(*_INK)
                        pdf.multi_cell(_PW, 6.2, line)
                        pdf.ln(1.5)

        pdf.ln(4)

    return bytes(pdf.output())


# ---------------------------------------------------------------------------
# PPTX generator
# ---------------------------------------------------------------------------

def _generate_pptx(project_id: str, run_id: str, sections: dict) -> bytes:
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    # ── Palette (matches PDF) ────────────────────────────────────────────────
    C_GREEN    = RGBColor(0x2C, 0x4A, 0x3E)
    C_OFFWHITE = RGBColor(0xF7, 0xF4, 0xEE)
    C_INK      = RGBColor(0x1A, 0x17, 0x13)
    C_INK_SOFT = RGBColor(0x50, 0x4B, 0x44)
    C_GRAY     = RGBColor(0x96, 0x90, 0x88)
    C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    C_DANGER   = RGBColor(0xC0, 0x39, 0x2B)
    C_WARNING  = RGBColor(0xD4, 0x8C, 0x10)
    C_SUCCESS  = RGBColor(0x27, 0x6E, 0x48)
    C_RULE     = RGBColor(0xD2, 0xCE, 0xC6)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    proj_display = project_id.replace("_", " ").replace("-", " ").title()
    date_str = datetime.now(timezone.utc).strftime("%B %d, %Y")

    # ── Low-level helpers ────────────────────────────────────────────────────

    def _rect(slide, x: float, y: float, w: float, h: float, color: RGBColor):
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def _tb(slide, x: float, y: float, w: float, h: float):
        return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

    def _p0(tb):
        """Return first paragraph of a freshly-created textbox."""
        tf = tb.text_frame
        tf.word_wrap = True
        return tf.paragraphs[0]

    def _set(para, text: str, size: float, bold=False, italic=False,
             color=C_INK, align=PP_ALIGN.LEFT):
        para.text = _safe(str(text))
        para.alignment = align
        if para.runs:
            r = para.runs[0]
        else:
            r = para.add_run()
            r.text = _safe(str(text))
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.italic = italic
        r.font.color.rgb = color

    def _bg(slide, color=C_OFFWHITE):
        _rect(slide, 0, 0, 13.33, 7.5, color)

    def _accent(slide):
        _rect(slide, 0, 0, 13.33, 0.07, C_GREEN)

    def _eyebrow(slide, text: str, x=0.5, y=0.18):
        tb = _tb(slide, x, y, 12, 0.3)
        _set(_p0(tb), text.upper(), 7.5, bold=True, color=C_GRAY)

    def _rule(slide, x=0.5, y=1.32, w=12.3):
        _rect(slide, x, y, w, 0.025, C_RULE)

    def _title(slide, text: str, y=0.55, size=30):
        tb = _tb(slide, 0.5, y, 12.3, 0.8)
        _set(_p0(tb), text, size, bold=True, color=C_INK)

    # ── Slide factory ────────────────────────────────────────────────────────

    def _content_slide(title: str, eyebrow: str):
        s = prs.slides.add_slide(blank)
        _bg(s)
        _accent(s)
        _eyebrow(s, eyebrow)
        _title(s, title)
        _rule(s)
        return s

    # ── Slide 1: Cover ───────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _bg(sl, C_GREEN)
    _accent(sl)

    tb = _tb(sl, 0.75, 1.6, 11.5, 1.6)
    tf = tb.text_frame
    tf.word_wrap = True
    _set(tf.paragraphs[0], proj_display, 44, bold=True, color=C_WHITE)

    assembly = sections.get("07_assembly", {})
    stage = ""
    if isinstance(assembly, dict):
        stage = str(assembly.get("project_stage") or assembly.get("study_level") or "")
    if stage and stage.lower() not in ("unknown", "not specified", ""):
        tb2 = _tb(sl, 0.75, 3.35, 11.0, 0.5)
        _set(_p0(tb2), stage, 17, color=RGBColor(0xA8, 0xC5, 0xB8))

    tb3 = _tb(sl, 0.75, 4.1, 11.0, 0.4)
    _set(_p0(tb3), f"{date_str}  ·  Run {run_id}  ·  Internal / Confidential", 10.5, color=C_GRAY)

    tb4 = _tb(sl, 0.75, 6.55, 11.5, 0.7)
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    _set(tf4.paragraphs[0],
         "For internal research purposes only. Not investment advice. "
         "Verify all figures against primary source documents.", 9,
         color=RGBColor(0x7A, 0x8A, 0x84))

    tb5 = _tb(sl, 11.3, 6.9, 1.8, 0.4)
    _set(_p0(tb5), "Extract", 11, bold=True, color=C_GRAY, align=PP_ALIGN.RIGHT)

    # ── Slide 2: Key Metrics ──────────────────────────────────────────────────
    callouts = []
    if isinstance(assembly, dict):
        raw = assembly.get("key_callouts")
        if isinstance(raw, list):
            callouts = raw

    if callouts:
        sl = _content_slide("At a Glance", "Key Metrics")
        cols = min(len(callouts), 3)
        rows = (min(len(callouts), 6) + cols - 1) // cols
        cw = 11.5 / cols
        ch = 4.9 / rows
        x0, y0 = 0.9, 1.5

        for i, c in enumerate(callouts[:6]):
            col, row = i % cols, i // cols
            cx = x0 + col * cw
            cy = y0 + row * ch
            _rect(sl, cx, cy, cw - 0.15, ch - 0.12, C_WHITE)

            tb_v = _tb(sl, cx + 0.2, cy + 0.15, cw - 0.5, 0.7)
            _set(_p0(tb_v), str(c.get("value", "")), 28, bold=True, color=C_INK)

            tb_l = _tb(sl, cx + 0.2, cy + 0.82, cw - 0.5, 0.35)
            _set(_p0(tb_l), str(c.get("label", "")), 11, color=C_INK_SOFT)

            ctx = c.get("context")
            if ctx:
                tb_c = _tb(sl, cx + 0.2, cy + 1.16, cw - 0.5, 0.32)
                _set(_p0(tb_c), str(ctx), 9.5, italic=True, color=C_GRAY)

    # ── Slide 3: Analyst Narrative ────────────────────────────────────────────
    narrative  = str(assembly.get("narrative") or "")  if isinstance(assembly, dict) else ""
    conclusion = str(assembly.get("analyst_conclusion") or "") if isinstance(assembly, dict) else ""

    if narrative:
        sl = _content_slide("Analyst Summary", "Analyst Narrative")
        y = 1.5
        for para in [p.strip() for p in narrative.split("\n\n") if p.strip()][:3]:
            if y > 5.6: break
            tb_p = _tb(sl, 0.5, y, 12.3, 1.05)
            _set(_p0(tb_p), para[:380], 13, color=C_INK)
            y += 1.1
        if conclusion:
            _rule(sl, y=5.85)
            tb_c = _tb(sl, 0.5, 5.95, 12.3, 1.2)
            tf_c = tb_c.text_frame
            tf_c.word_wrap = True
            _set(tf_c.paragraphs[0], conclusion[:280], 12, italic=True, color=C_INK_SOFT)

    # ── Slides 4–6: Specialist sections ──────────────────────────────────────

    def _kv_from(d: dict, text_keys: list, kv_keys: list) -> tuple[str, list]:
        prose = ""
        for k in text_keys:
            v = d.get(k)
            if isinstance(v, str) and len(v) > 40:
                prose = v
                break
        kv = [(k, v) for k in kv_keys
              if (v := d.get(k)) is not None and not isinstance(v, (dict, list))]
        return prose, kv

    def _draw_kv_slide(sl, kv: list, prose: str):
        """Render KV grid or prose body on an existing content slide."""
        y = 1.5
        if prose and not kv:
            for para in [p.strip() for p in prose.split("\n\n") if p.strip()][:3]:
                if y > 6.5: break
                tb_p = _tb(sl, 0.5, y, 12.3, 1.0)
                _set(_p0(tb_p), para[:380], 13, color=C_INK)
                y += 1.1
            return
        # Two-column KV grid
        left_col  = kv[0::2]
        right_col = kv[1::2]
        for col_idx, col in enumerate((left_col, right_col)):
            x = 0.5 + col_idx * 6.25
            cy = y
            for lbl, val in col[:7]:
                if cy > 6.6: break
                tb_lbl = _tb(sl, x, cy, 5.8, 0.27)
                _set(_p0(tb_lbl), str(lbl).replace("_", " ").title(), 8.5, bold=True, color=C_GRAY)
                tb_val = _tb(sl, x, cy + 0.26, 5.8, 0.38)
                tf_val = tb_val.text_frame
                tf_val.word_wrap = True
                _set(tf_val.paragraphs[0], str(val)[:120], 12.5, color=C_INK)
                cy += 0.75

    geo = sections.get("03_geology", {})
    if isinstance(geo, dict) and geo:
        prose, kv = _kv_from(geo,
            ["geology_summary", "deposit_description", "summary", "geology", "overview"],
            ["deposit_type", "commodity", "jurisdiction", "resource_estimate",
             "total_resource_mt", "grade", "contained_metal", "classification"])
        sl = _content_slide("Geology & Resources", "Section 1")
        _draw_kv_slide(sl, kv, prose)

    econ = sections.get("04_economics", {})
    if isinstance(econ, dict) and econ:
        prose, kv = _kv_from(econ,
            ["economics_summary", "summary", "overview", "economics"],
            ["npv_musd", "irr_percent", "payback_years", "initial_capex_musd",
             "aisc_per_oz", "opex_per_tonne", "mine_life_years",
             "commodity_price_assumption", "discount_rate"])
        sl = _content_slide("Economics & Financial Analysis", "Section 2")
        _draw_kv_slide(sl, kv, prose)

    risks = sections.get("05_risks", {})
    if isinstance(risks, dict) and risks:
        risk_list  = risks.get("risks") or risks.get("key_risks") or risks.get("risk_factors") or []
        risk_prose = str(risks.get("summary") or risks.get("risk_overview") or "")
        sl = _content_slide("Risks & Uncertainties", "Section 3")
        y = 1.5
        if isinstance(risk_list, list) and risk_list:
            for risk in risk_list[:7]:
                if y > 6.6: break
                lbl_t = desc_t = ""
                if isinstance(risk, dict):
                    lbl_t  = str(risk.get("risk_type") or risk.get("category") or risk.get("risk") or "")
                    desc_t = str(risk.get("description") or risk.get("detail") or risk.get("mitigation") or "")
                else:
                    desc_t = str(risk)
                _rect(sl, 0.5, y + 0.14, 0.06, 0.06, C_GREEN)
                tb_r = _tb(sl, 0.72, y, 12.0, 0.7)
                tf_r = tb_r.text_frame
                tf_r.word_wrap = True
                p_r = tf_r.paragraphs[0]
                if lbl_t:
                    _set(p_r, lbl_t, 12, bold=True, color=C_INK)
                    if desc_t:
                        p2 = tf_r.add_paragraph()
                        _set(p2, desc_t[:200], 11, color=C_INK_SOFT)
                else:
                    _set(p_r, desc_t[:230], 12, color=C_INK)
                y += 0.82
        elif risk_prose:
            tb_p = _tb(sl, 0.5, y, 12.3, 4.5)
            tf_p = tb_p.text_frame
            tf_p.word_wrap = True
            _set(tf_p.paragraphs[0], risk_prose[:600], 13, color=C_INK)

    # ── Slide 7: DCF summary ──────────────────────────────────────────────────
    dcf = sections.get("06_dcf_model", {})
    if isinstance(dcf, dict) and dcf.get("model_ran"):
        summary = dcf.get("summary") or {}
        if isinstance(summary, dict):
            skip = {"project_id", "scenario", "after_tax", "notes", "aisc_unit"}
            kv = [(k, f"{v:,}" if isinstance(v, (int, float)) else str(v))
                  for k, v in summary.items()
                  if k not in skip and v is not None]
            sl = _content_slide("DCF Financial Model", "Section 4")
            _draw_kv_slide(sl, kv, "")

    # ── Slide 8: Data Quality ─────────────────────────────────────────────────
    gaps = sections.get("08_data_gaps", {})
    conf = sections.get("09_confidence", {})
    if (isinstance(gaps, dict) and gaps) or (isinstance(conf, dict) and conf):
        sl = _content_slide("Data Quality Assessment", "Data Quality")

        # Left column — gaps
        if isinstance(gaps, dict) and gaps:
            eyebrow_l = _tb(sl, 0.5, 1.5, 5.9, 0.3)
            _set(_p0(eyebrow_l), "DATA GAPS", 8.5, bold=True, color=C_GRAY)

            counts = []
            if gaps.get("critical_gaps_count"):  counts.append(f"{gaps['critical_gaps_count']} critical")
            if gaps.get("important_gaps_count"): counts.append(f"{gaps['important_gaps_count']} important")
            if gaps.get("minor_gaps_count"):     counts.append(f"{gaps['minor_gaps_count']} minor")
            if counts:
                tb_cnt = _tb(sl, 0.5, 1.85, 5.9, 0.45)
                _set(_p0(tb_cnt), "  ·  ".join(counts), 14, bold=True, color=C_INK)

            overall_g = str(gaps.get("overall_data_quality_comment") or "")
            if overall_g:
                tb_og = _tb(sl, 0.5, 2.4, 5.9, 2.2)
                tf_og = tb_og.text_frame
                tf_og.word_wrap = True
                _set(tf_og.paragraphs[0], overall_g[:320], 12, color=C_INK_SOFT)

        # Divider
        _rect(sl, 6.66, 1.45, 0.03, 5.6, C_RULE)

        # Right column — confidence
        if isinstance(conf, dict) and conf:
            eyebrow_r = _tb(sl, 6.83, 1.5, 5.9, 0.3)
            _set(_p0(eyebrow_r), "CONFIDENCE", 8.5, bold=True, color=C_GRAY)

            overall_c = str(conf.get("overall_confidence_statement") or "")
            if overall_c:
                tb_oc = _tb(sl, 6.83, 1.85, 5.9, 2.0)
                tf_oc = tb_oc.text_frame
                tf_oc.word_wrap = True
                _set(tf_oc.paragraphs[0], overall_c[:300], 12, color=C_INK_SOFT)

            y_q = 4.0
            for lbl_str, val_str, bar_color in [
                ("Most reliable",  str(conf.get("most_reliable_aspect")  or ""), C_SUCCESS),
                ("Least reliable", str(conf.get("least_reliable_aspect") or ""), C_WARNING),
            ]:
                if not val_str: continue
                _rect(sl, 6.83, y_q, 0.05, 0.6, bar_color)
                tb_q = _tb(sl, 7.02, y_q, 5.9, 0.65)
                tf_q = tb_q.text_frame
                tf_q.word_wrap = True
                _set(tf_q.paragraphs[0], f"{lbl_str}: {val_str[:130]}", 11, color=C_INK)
                y_q += 0.82

    # ── Slide 9: Contradictions ───────────────────────────────────────────────
    contras = sections.get("10_contradictions", {})
    if isinstance(contras, dict):
        total_c = (contras.get("contradiction_count") or 0) + len(contras.get("arithmetic_errors") or [])
        if total_c:
            sl = _content_slide("Contradiction & Consistency Flags", "QA Review")
            crit_n = contras.get("critical_count",    0) or 0
            sig_n  = contras.get("significant_count", 0) or 0
            min_n  = contras.get("minor_count",       0) or 0

            # Chip row
            chip_x = 0.5
            for chip_lbl, chip_val, chip_col in [
                ("CRITICAL",    crit_n, C_DANGER),
                ("SIGNIFICANT", sig_n,  C_WARNING),
                ("MINOR",       min_n,  C_GRAY),
            ]:
                _rect(sl, chip_x, 1.5, 1.85, 0.75, chip_col)
                tb_cv = _tb(sl, chip_x + 0.05, 1.52, 1.75, 0.42)
                _set(_p0(tb_cv), str(chip_val), 22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
                tb_cl2 = _tb(sl, chip_x + 0.05, 1.93, 1.75, 0.28)
                _set(_p0(tb_cl2), chip_lbl, 7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
                chip_x += 2.0

            # Top 5 contradictions
            all_c = sorted(
                [c for c in (contras.get("contradictions") or []) if isinstance(c, dict)],
                key=lambda x: {"critical": 0, "significant": 1, "minor": 2}.get(x.get("severity", "minor"), 2),
            )
            y = 2.5
            for c in all_c[:5]:
                if y > 6.6: break
                sev = c.get("severity", "minor")
                bar_col = {"critical": C_DANGER, "significant": C_WARNING}.get(sev, C_GRAY)
                _rect(sl, 0.5, y, 0.05, 0.58, bar_col)
                tb_ci = _tb(sl, 0.72, y, 12.1, 0.65)
                tf_ci = tb_ci.text_frame
                tf_ci.word_wrap = True
                _set(tf_ci.paragraphs[0], str(c.get("description", ""))[:220], 12, color=C_INK)
                y += 0.78

    # ── NI 43-101 / JORC Compliance slide ────────────────────────────────────
    compliance = sections.get("13_compliance", {})
    if isinstance(compliance, dict) and compliance.get("overall_status"):
        sl = _content_slide("NI 43-101 / JORC Compliance", "Reporting Standards Review")
        overall = compliance.get("overall_status", "")
        std_applied = compliance.get("standard_applied_for_assessment", "NI 43-101")
        summary_txt = compliance.get("overall_summary", "")
        met_n     = int(compliance.get("compliant_items_count", 0) or 0)
        partial_n = int(compliance.get("partial_items_count",   0) or 0)
        missing_n = int(compliance.get("missing_items_count",   0) or 0)
        crit_gaps = compliance.get("critical_gaps") or []
        study_type = compliance.get("study_type", "")
        study_match = compliance.get("study_type_resource_match", "")
        qp_data = compliance.get("qualified_person") or {}

        _OVERALL_COL = {
            "compliant":          C_GREEN,
            "likely_compliant":   RGBColor(0x1e, 0x40, 0xaf),
            "deficiencies_found": C_WARNING,
            "major_gaps":         C_DANGER,
        }
        _OVERALL_LBL = {
            "compliant":          "COMPLIANT",
            "likely_compliant":   "LIKELY COMPLIANT",
            "deficiencies_found": "DEFICIENCIES FOUND",
            "major_gaps":         "MAJOR GAPS",
        }
        ov_col = _OVERALL_COL.get(overall, C_GRAY)

        # Overall status banner
        _rect(sl, 0.5, 1.5, 5.5, 0.72, ov_col)
        tb_ov = _tb(sl, 0.55, 1.52, 5.4, 0.68)
        _set(_p0(tb_ov), _OVERALL_LBL.get(overall, overall.upper()), 14,
             bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Standard badge
        tb_std = _tb(sl, 6.2, 1.52, 6.3, 0.52)
        _set(_p0(tb_std), f"Standard: {std_applied}", 11, color=C_INK)

        # Count chips
        chip_x = 0.5
        for lbl, val, col in [("MET", met_n, C_GREEN), ("PARTIAL", partial_n, C_WARNING), ("MISSING", missing_n, C_DANGER)]:
            _rect(sl, chip_x, 2.42, 1.75, 0.65, col)
            tb_v = _tb(sl, chip_x + 0.05, 2.44, 1.65, 0.36)
            _set(_p0(tb_v), str(val), 20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            tb_l = _tb(sl, chip_x + 0.05, 2.79, 1.65, 0.24)
            _set(_p0(tb_l), lbl, 7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            chip_x += 1.9

        # QP + study type row
        y_info = 3.25
        if isinstance(qp_data, dict):
            qp_name = qp_data.get("name") or ""
            qp_named = bool(qp_data.get("named"))
            qp_line = (f"QP: {qp_name}" if qp_named and qp_name else "QP: Not Named")
            tb_qp = _tb(sl, 0.5, y_info, 5.8, 0.42)
            _set(_p0(tb_qp), qp_line, 11, color=C_INK)

        if study_type:
            match_label = {"ok": "Study/Resource: Consistent", "concern": "Study/Resource: Concern", "violation": "Study/Resource: VIOLATION"}.get(study_match, "")
            sm_col = {"ok": C_GREEN, "concern": C_WARNING, "violation": C_DANGER}.get(study_match, C_GRAY)
            tb_sm = _tb(sl, 6.5, y_info, 5.8, 0.42)
            _set(_p0(tb_sm), f"{study_type.upper()}  —  {match_label}" if match_label else study_type.upper(), 11, color=sm_col)

        # Summary
        if summary_txt:
            tb_s = _tb(sl, 0.5, 3.8, 12.0, 0.85)
            tf_s = tb_s.text_frame
            tf_s.word_wrap = True
            _set(tf_s.paragraphs[0], str(summary_txt)[:300], 10.5, color=C_INK)

        # Critical gaps list
        if crit_gaps:
            tb_gh = _tb(sl, 0.5, 4.8, 2.5, 0.3)
            _set(_p0(tb_gh), "CRITICAL GAPS", 8.5, bold=True, color=C_DANGER)
            y_gap = 5.2
            for gap in crit_gaps[:4]:
                if y_gap > 6.6: break
                _rect(sl, 0.5, y_gap, 0.04, 0.4, C_DANGER)
                tb_g = _tb(sl, 0.68, y_gap, 11.8, 0.45)
                tf_g = tb_g.text_frame
                tf_g.word_wrap = True
                _set(tf_g.paragraphs[0], str(gap)[:200], 10.5, color=C_INK)
                y_gap += 0.55

    # ── Closing slide ─────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _bg(sl, C_GREEN)
    _accent(sl)

    tb = _tb(sl, 0.75, 2.4, 11.0, 1.3)
    tf = tb.text_frame
    tf.word_wrap = True
    _set(tf.paragraphs[0], proj_display, 38, bold=True, color=C_WHITE)

    tb2 = _tb(sl, 0.75, 3.85, 10.5, 0.5)
    _set(_p0(tb2), "Technical Analysis — Internal Use Only", 14,
         color=RGBColor(0xA8, 0xC5, 0xB8))

    tb3 = _tb(sl, 0.75, 6.8, 10.5, 0.4)
    _set(_p0(tb3), f"Extract  ·  {date_str}", 10, color=C_GRAY)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Route
# ---------------------------------------------------------------------------

@router.get("/projects/{project_id}/reports/{run_id}/export")
def export_report(
    project_id: str,
    run_id: str,
    format: Literal["json", "md", "txt", "pdf", "pptx"] = Query("pdf", description="Export format"),
) -> Response:
    _project_exists(project_id)
    sections = _get_all_sections(project_id, run_id)
    if not sections:
        raise HTTPException(status_code=404, detail="No report output found for this run.")

    filename_base = f"{project_id}_{run_id}_report"

    if format == "json":
        payload = json.dumps({"project_id": project_id, "run_id": run_id, "sections": sections}, indent=2)
        return Response(
            content=payload,
            media_type="application/json",
            headers={"Content-Disposition": f'attachment; filename="{filename_base}.json"'},
        )

    if format == "pdf":
        pdf_bytes = _generate_pdf(project_id, run_id, sections)
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename_base}.pdf"'},
        )

    if format == "pptx":
        pptx_bytes = _generate_pptx(project_id, run_id, sections)
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename_base}.pptx"'},
        )

    md_content = _sections_to_markdown(project_id, run_id, sections)

    if format == "txt":
        txt = re.sub(r"#{1,6} ", "", md_content)
        txt = re.sub(r"\*\*(.+?)\*\*", r"\1", txt)
        txt = re.sub(r"_(.+?)_", r"\1", txt)
        return Response(
            content=txt,
            media_type="text/plain",
            headers={"Content-Disposition": f'attachment; filename="{filename_base}.txt"'},
        )

    return Response(
        content=md_content,
        media_type="text/markdown",
        headers={"Content-Disposition": f'attachment; filename="{filename_base}.md"'},
    )
